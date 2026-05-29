#!/usr/bin/env python3
"""
Generador de presentación Session 5: A.8, Certificación ISO 27001, y Continuidad de Negocio
Requiere: python-pptx
Ejecutar: python generate_presentation.py
"""

import subprocess
import sys
from pathlib import Path

# Auto-instalar python-pptx si no está disponible
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Instalando python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

# Colores
D = RGBColor(0x1C, 0x1C, 0x1C)
W = RGBColor(0xFF, 0xFF, 0xFF)
G = RGBColor(0xF3, 0xF3, 0xF3)
DG = RGBColor(0x2D, 0x2D, 0x2D)
MUT = RGBColor(0x6B, 0x6B, 0x6B)
SMUT = RGBColor(0x9A, 0x9A, 0x9A)
WMUT = RGBColor(0x88, 0x88, 0x88)
ACC = RGBColor(0xD6, 0x4D, 0x00)
BLU = RGBColor(0x1B, 0x5F, 0xA8)
GRN = RGBColor(0x1A, 0x5C, 0x35)
TEA = RGBColor(0x0E, 0x66, 0x74)
RED = RGBColor(0x8B, 0x1A, 0x1A)

# Crear presentación
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """Añade un rectángulo a la slide"""
    shape = slide.shapes.add_shape(
        1,  # msoShapeRectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape

def add_text(slide, text, left, top, width, height, font_size=14, color=W, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """Añade texto a la slide"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def dark_slide():
    """Crea una slide con fondo oscuro"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = D
    return slide

def white_slide():
    """Crea una slide con fondo blanco"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = W
    return slide

def section_slide(title, subtitle=""):
    """Crea una slide de sección"""
    slide = dark_slide()
    add_text(slide, title, 0.65, 1.75, 9, 0.9, font_size=50, color=W, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.65, 2.78, 9, 0.42, font_size=17, color=WMUT)
    return slide

# ─────────────────────────────────────────
# SLIDE 01: COVER
# ─────────────────────────────────────────
slide = dark_slide()
add_text(slide, "Sesión 5", 0.65, 0.7, 8, 0.35, font_size=15, color=WMUT)
add_text(slide, "A.8 · Certificación ISO 27001\nDeclaración ENS\n+ Intro Continuidad de Negocio",
         0.65, 1.08, 9, 2.5, font_size=42, color=W, bold=True)
add_text(slide, "Retomamos S4 · Avanzamos hacia el módulo de BCM · ISO 22301 · ENS [op.cont.*]",
         0.65, 3.7, 9, 0.35, font_size=13, color=WMUT)
add_shape(slide, 0.65, 4.2, 1.1, 0.04, WMUT)
add_text(slide, "ISO 27001:2022 + ENS (RD 311/2022)", 0.65, 4.38, 8, 0.28, font_size=11, color=SMUT)

# ─────────────────────────────────────────
# SLIDE 02: AGENDA
# ─────────────────────────────────────────
slide = dark_slide()
add_text(slide, "Agenda", 0.65, 0.55, 4, 0.58, font_size=36, color=W, bold=True)
add_text(slide, "Sesión 5 · 3 horas", 0.65, 1.18, 4, 0.3, font_size=14, color=WMUT)

items = [
    ("00", "Repaso S4 · Kahoot 10 preguntas", "0:00–0:25", ""),
    ("01", "[Pendiente S4] Controles A.8 — Los 34 controles tecnológicos", "0:25–1:05", "S4 B2"),
    ("02", "[Pendiente S4] Certificación ISO 27001 · Declaración ENS", "1:05–1:25", "S4 B3"),
    ("⏸", "PAUSA", "1:25–1:45", ""),
    ("03", "Continuidad de Negocio — Introducción y fundamentos", "1:45–2:50", "BCM"),
    ("✓", "Síntesis · Preview S6 (BCM completo)", "2:50–3:00", ""),
]

for i, (n, t, tm, note) in enumerate(items):
    y = 1.6 + i * 0.64
    color = ACC if n == "⏸" else TEA if note else WMUT
    add_shape(slide, 5.1, y, 4.6, 0.56, DG)
    add_text(slide, n, 5.1, y, 0.55, 0.56, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, t, 5.7, y + 0.06, 3.85, 0.26, font_size=12.5, color=W, bold=True)
    add_text(slide, tm, 5.7, y + 0.34, 3.85, 0.2, font_size=10, color=SMUT)

# ─────────────────────────────────────────
# SLIDE 03: A.8 INTRO
# ─────────────────────────────────────────
section_slide("Bloque 01 [pendiente S4]\nControles A.8", "Los 34 controles tecnológicos · Comparativa ISO 27001 ↔ ENS")

# ─────────────────────────────────────────
# SLIDE 04: A.8 CONTROLS I
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "A.8 · Controles tecnológicos I", 0.38, 0.16, 9.2, 0.44, font_size=20, color=D, bold=True)
add_text(slide, "Endpoints, accesos privilegiados, malware y vulnerabilidades — A.8.1 a A.8.8",
         0.38, 0.63, 9.2, 0.27, font_size=11, color=MUT)

# Header
headers = ["Control", "Propósito y riesgo", "Evidencia clave", "ENS equivalente (825)"]
col_x = [0.38, 2.05, 4.22, 6.4]
col_w = [1.62, 2.12, 2.12, 2.88]

for i, h in enumerate(headers):
    add_shape(slide, col_x[i], 0.95, col_w[i] - 0.02, 0.3, D)
    add_text(slide, h, col_x[i] + 0.06, 0.95, col_w[i] - 0.1, 0.3,
             font_size=9.5, color=W, bold=True, valign=MSO_ANCHOR.MIDDLE)

# Datos de controles A.8.1 a A.8.8
controls_1 = [
    ("A.8.1", "Dispositivos usuario final", "Datos corporativos en dispositivos no gestionados",
     "Política, inventario, MDM con cifrado", "[mp.eq.1]", "Puesto de trabajo"),
    ("A.8.2", "Accesos privilegiados", "Cuenta privilegiada comprometida → acceso total",
     "Inventario, PAM, auditoría", "[op.acc.4]", "Gestión derechos privilegiados"),
    ("A.8.5", "Autenticación segura", "Credenciales robadas o débiles",
     "MFA, política de contraseñas en AD", "[op.acc.6]", "MFA requerido"),
]

for i, (code, name, risk, ev, ens, ensd) in enumerate(controls_1):
    y = 1.28 + i * 0.53
    bg = W if i % 2 == 0 else G
    add_shape(slide, 0.38, y, 9.2, 0.51, bg)
    add_shape(slide, 0.38, y, 0.05, 0.51, TEA)
    add_text(slide, code, 0.48, y + 0.04, 0.72, 0.2, font_size=8.5, color=TEA, bold=True)
    add_text(slide, name, 0.48, y + 0.25, 1.52, 0.22, font_size=10, color=D, bold=True)
    add_text(slide, risk, col_x[1] + 0.05, y + 0.04, col_w[1] - 0.12, 0.44, font_size=10, color=D)
    add_text(slide, ev, col_x[2] + 0.05, y + 0.04, col_w[2] - 0.12, 0.44, font_size=10, color=MUT)
    add_shape(slide, col_x[3] + 0.05, y + 0.08, col_w[3] - 0.15, 0.26, GRN)
    add_text(slide, ens, col_x[3] + 0.1, y + 0.08, col_w[3] - 0.22, 0.26, font_size=9, color=W, bold=True, valign=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────
# SLIDE 05: EJERCICIO A.8
# ─────────────────────────────────────────
slide = dark_slide()
add_text(slide, "Ejercicio · 8 minutos", 0.65, 0.5, 8, 0.28, font_size=12, color=SMUT, bold=True)
add_text(slide, "TechServ: los 5 controles A.8 más urgentes", 0.65, 0.85, 9, 0.58, font_size=34, color=W, bold=True)
add_text(slide, "Sistemas ALTA: nóminas · expedientes · sede electrónica. Sin SGSI previo.",
         0.65, 1.5, 9, 0.32, font_size=13, color=WMUT)
add_shape(slide, 0.5, 1.95, 9.1, 2.65, DG)
add_text(slide, "Para cada control elegido indicad: código · riesgo · acción · evidencia · ENS equivalente",
         0.7, 2.05, 8.7, 0.35, font_size=12, color=W, bold=True)

hints = [
    "A.8.5 MFA — accesos remotos de personal y usuarios de 3 Consejerías",
    "A.8.8 Vulnerabilidades — sistemas críticos para AAPP",
    "A.8.13 Backup — pérdida = catástrofe legal",
    "A.8.12 DLP — datos de ciudadanos y funcionarios",
    "A.8.16 Monitorización — detectar accesos anómalos",
    "A.8.22 Segmentación — separar entornos críticos",
]

for i, hint in enumerate(hints):
    add_text(slide, "→ " + hint, 0.88, 2.5 + i * 0.35, 8.3, 0.32, font_size=11.5, color=WMUT)

add_shape(slide, 0.5, 4.68, 9.1, 0.7, D)
add_text(slide, "Puesta en común: 1 portavoz · debate sobre exclusiones justificadas",
         0.7, 4.72, 8.7, 0.6, font_size=11, color=WMUT)

# ─────────────────────────────────────────
# SLIDE 06: CERTIFICACIÓN INTRO
# ─────────────────────────────────────────
section_slide("Bloque 02 [pendiente S4]\nCertificación ISO 27001\ny Declaración ENS",
              "El proceso completo · Qué evalúa el auditor · Diferencias entre ambos")

# ─────────────────────────────────────────
# SLIDE 07: CERTIFICACIÓN 7 FASES
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "El camino a la certificación ISO 27001", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "7 fases · Tiempos orientativos · Lo que evalúa el auditor en cada paso",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

steps = [
    ("1", "GAP Analysis", "Diagnóstico estado actual", "2-4 sem", DG),
    ("2", "Implementación", "SoA, PTR, políticas, controles", "6-18 m", DG),
    ("3", "Auditoría interna", "Cl. 9.2: auditar SGSI propio", "1-2 sem", "333333"),
    ("4", "Revisión dirección", "Cl. 9.3: aprobación. Acta.", "1 día", DG),
    ("5", "Etapa 1 Doc.", "Auditor revisa documentación", "1-2 d", "333333"),
    ("6", "Etapa 2 In situ", "Verifica implementación real", "2-5 d", DG),
    ("7", "Certificado 3 años", "Emitido ENAC. Seguimiento anual", "3 años", GRN),
]

for i, (n, t, d, tm, c) in enumerate(steps):
    x = 0.38 + i * 1.35
    add_shape(slide, x, 1.05, 1.22, 0.55, RGBColor(int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)) if isinstance(c, str) and c.startswith('#') else c)
    add_text(slide, n, x + 0.04, 1.08, 0.3, 0.48, font_size=19, color=WMUT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, t, x + 0.38, 1.1, 0.82, 0.45, font_size=9.5, color=W, bold=True)

    add_shape(slide, x, 1.68, 1.22, 3.1, G if i % 2 == 0 else W)
    add_text(slide, tm, x + 0.04, 1.7, 1.15, 0.21, font_size=8.5, color=W, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, d, x + 0.06, 1.98, 1.12, 2.65, font_size=10, color=D)

add_shape(slide, 0.38, 4.88, 9.22, 0.58, D)
add_text(slide, "TechServ sin SGSI previo: 12–18 meses. Con base: 6–12 meses. Etapa 1 falla → implementación. Etapa 2 NC mayor → espera",
         0.58, 4.92, 8.95, 0.5, font_size=11.5, color=WMUT)

# ─────────────────────────────────────────
# SLIDE 08: DECLARACIÓN ENS
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "Declaración de Conformidad ENS", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "Proceso · Portal CCN · Tabla comparativa ISO 27001 vs. Declaración ENS",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

# Left side: Proceso
add_shape(slide, 0.45, 0.95, 4.55, 4.42, G)
add_text(slide, "Proceso — Declaración ENS", 0.62, 1.04, 4.22, 0.28, font_size=13, color=D, bold=True)

process = [
    "Categorizar sistema según Anexo I",
    "Aplicar medidas Anexo II por categoría",
    "Auditoría ENS: cada 2 años (MEDIA/ALTA)",
    "Organismo emite: conforme / no conforme",
    "Entidad publica en: ens.ccn.cni.es"
]

for i, p in enumerate(process):
    y = 1.42 + i * 0.72
    add_shape(slide, 0.62, y, 0.44, 0.54, D)
    add_text(slide, str(i + 1), 0.62, y, 0.44, 0.54, font_size=16, color=W, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_shape(slide, 1.1, y, 3.72, 0.54, W)
    add_text(slide, p, 1.2, 1.5 + i * 0.72, 3.56, 0.38, font_size=11, color=D)

# Right side: Comparativa
add_shape(slide, 5.18, 0.95, 4.42, 4.42, D)
add_text(slide, "ISO 27001 vs. Declaración ENS", 5.36, 1.04, 4.08, 0.28, font_size=13, color=W, bold=True)

comparision = [
    ("Naturaleza", "Certificación externa", "Declaración + auditoría"),
    ("¿Quién certifica?", "ENAC", "La entidad + auditor ENS"),
    ("Periodicidad", "Anual + 3 años", "Cada 2 años (MEDIA/ALTA)"),
    ("Reconocimiento", "Internacional", "Nacional (España)"),
    ("Ámbito", "Toda org. o parte", "Por sistema de información"),
    ("Catálogo público", "No hay registro único", "Portal ENS del CCN"),
]

for i, (d, a, b) in enumerate(comparision):
    y = 1.42 + i * 0.62
    add_shape(slide, 5.3, y, 1.6, 0.54, DG)
    add_text(slide, d, 5.38, y + 0.09, 1.48, 0.36, font_size=9.5, color=WMUT)
    add_shape(slide, 6.95, y, 1.32, 0.54, DG)
    add_text(slide, a, 7.0, y + 0.09, 1.22, 0.36, font_size=9.5, color=BLU, bold=True)
    add_shape(slide, 8.32, y, 1.22, 0.54, DG)
    add_text(slide, b, 8.37, y + 0.09, 1.12, 0.36, font_size=9.5, color=GRN, bold=True)

add_text(slide, "ISO", 6.95, 1.04, 1.32, 0.28, font_size=10, color=W, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_text(slide, "ENS", 8.32, 1.04, 1.22, 0.28, font_size=10, color=W, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────
# SLIDE 09: PAUSA
# ─────────────────────────────────────────
slide = dark_slide()
add_text(slide, "⏸", 4.0, 1.4, 2, 1.5, font_size=90, color=WMUT, align=PP_ALIGN.CENTER)
add_text(slide, "Pausa", 2.5, 2.9, 5, 0.85, font_size=60, color=W, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, "20 minutos · Volvemos a las :", 2.5, 3.85, 5, 0.45, font_size=18, color=WMUT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 10: BCM INTRO
# ─────────────────────────────────────────
section_slide("Bloque 03\nContinuidad de Negocio", "BIA · RTO/RPO · BCP vs DRP · ENS [op.cont.*] · ISO 22301")

# ─────────────────────────────────────────
# SLIDE 11: ¿POR QUÉ BCM?
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "¿Por qué existe la continuidad de negocio?", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "Porque los incidentes ocurren. La pregunta no es si — es cuándo y qué tan preparados estáis.",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

cases = [
    ("Hospital Clínic · Barcelona · 2023", "RansomHouse cifra sistemas. 150 operaciones canceladas.",
     "Semanas sin servicio", RED),
    ("Ayuntamiento Sevilla · 2023", "LockBit3 ataca sistemas municipales. Servicios paralizados meses.",
     "Meses recuperación", RED),
    ("Maersk · Global · 2017", "NotPetya: 45.000 PCs destruidos. Restaurados en 10 días gracias a BCP.",
     "300M$ · 10 días", GRN),
]

for i, (org, d, cost, color) in enumerate(cases):
    y = 1.05 + i * 1.46
    add_shape(slide, 0.45, y, 9.1, 1.32, G if i % 2 == 0 else W)
    add_shape(slide, 0.45, y, 0.06, 1.32, color)
    add_text(slide, org, 0.6, y + 0.1, 3.5, 0.28, font_size=11, color=D, bold=True)
    add_text(slide, d, 0.6, y + 0.42, 6.0, 0.78, font_size=11.5, color=D)
    add_shape(slide, 6.72, y + 0.1, 2.7, 0.55, color)
    add_text(slide, cost, 6.82, y + 0.12, 2.55, 0.5, font_size=10, color=W)

# ─────────────────────────────────────────
# SLIDE 12: BIA
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "Business Impact Analysis (BIA)", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "El primer paso de cualquier plan de continuidad · ENS [op.cont.1]",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

# Left: What is BIA
add_shape(slide, 0.45, 0.95, 4.55, 4.6, G)
add_text(slide, "¿Qué es el BIA?", 0.62, 1.02, 4.22, 0.3, font_size=13, color=D, bold=True)
add_text(slide, "Análisis que identifica:\n① Funciones críticas\n② Impacto de interrupción\n③ Dependencias\n④ RTO/RPO requeridos\n\nSin BIA, el BCP es ficción.",
         0.62, 1.38, 4.22, 3.95, font_size=12, color=D)

# Right: BIA table
add_text(slide, "Ejemplo de BIA para TechServ", 5.18, 0.95, 4.45, 0.3, font_size=13, color=D, bold=True)

# Tabla simplificada
bia_data = [
    ("Sede electrónica", "ALTO", "MUY ALTO", "4 h", "1 h"),
    ("BBDD Nóminas", "MEDIO", "ALTO", "8 h", "0 h"),
    ("Expedientes", "BAJO", "ALTO", "24 h", "4 h"),
]

for i, (func, imp1, imp24, rto, rpo) in enumerate(bia_data):
    y = 1.62 + i * 0.72
    bg = W if i % 2 == 0 else G
    add_shape(slide, 5.18, y, 1.55, 0.68, bg)
    add_text(slide, func, 5.22, y + 0.06, 1.5, 0.56, font_size=10.5, color=D, bold=True)
    add_shape(slide, 6.76, y, 0.85, 0.68, bg)
    add_text(slide, imp1, 6.8, y + 0.06, 0.8, 0.56, font_size=9.5, color=MUT)
    add_shape(slide, 7.64, y, 0.85, 0.68, bg)
    add_text(slide, imp24, 7.68, y + 0.06, 0.8, 0.56, font_size=9.5, color=MUT)
    add_shape(slide, 8.52, y, 0.7, 0.68, bg)
    add_text(slide, rto, 8.56, y + 0.06, 0.65, 0.56, font_size=9.5, color=MUT)
    add_shape(slide, 9.25, y, 0.7, 0.68, bg)
    add_text(slide, rpo, 9.29, y + 0.06, 0.65, 0.56, font_size=9.5, color=MUT)

# ─────────────────────────────────────────
# SLIDE 13: RTO/RPO
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "RTO y RPO — Los dos parámetros clave", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "Recovery Time Objective · Recovery Point Objective · MAD / MTPD",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

# Timeline
add_shape(slide, 0.5, 1.05, 8.8, 0.06, D)
add_text(slide, "TIEMPO →", 9.0, 1.0, 0.7, 0.2, font_size=9, color=MUT)

# Key points
add_shape(slide, 2.5, 0.65, 0.05, 0.85, RED)
add_text(slide, "Incidente", 2.1, 0.42, 0.9, 0.38, font_size=10, color=RED, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, 1.2, 0.65, 0.05, 0.85, WMUT)
add_text(slide, "Backup", 0.8, 0.42, 0.8, 0.38, font_size=10, color=MUT, align=PP_ALIGN.CENTER)

add_shape(slide, 7.0, 0.65, 0.05, 0.85, GRN)
add_text(slide, "Restaurado", 6.7, 0.42, 0.9, 0.38, font_size=10, color=GRN, bold=True, align=PP_ALIGN.CENTER)

# Definitions
defs = [
    ("RTO — Recovery Time Objective", "Tiempo máximo para restaurar. Lo define el negocio. RTO bajo = mayor inversión.", BLU),
    ("RPO — Recovery Point Objective", "Datos que se pueden perder. RPO=0h → síncrono. RPO=24h → diario. Bajo = costoso.", RED),
    ("MAD / MTPD", "Tiempo máximo tolerable sin la función. Daño irreversible después. RTO < MAD siempre.", ACC),
]

for i, (t, d, color) in enumerate(defs):
    y = 1.75 + i * 1.2
    add_shape(slide, 0.45, y, 9.1, 1.12, G if i % 2 == 0 else W)
    add_shape(slide, 0.45, y, 0.07, 1.12, color)
    add_text(slide, t, 0.62, y + 0.02, 3.0, 0.28, font_size=12, color=D, bold=True)
    add_text(slide, d, 0.62, y + 0.32, 8.7, 0.64, font_size=11.5, color=MUT)

# ─────────────────────────────────────────
# SLIDE 14: BCP vs DRP
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "BCP y DRP — Dos planes complementarios", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "Business Continuity Plan · Disaster Recovery Plan",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

# Left: BCP
add_shape(slide, 0.45, 0.95, 4.5, 4.5, G)
add_shape(slide, 0.45, 0.95, 4.5, 0.06, BLU)
add_text(slide, "BCP — Plan de Continuidad de Negocio", 0.62, 1.06, 4.18, 0.28, font_size=13, color=D, bold=True)

bcp_points = [
    "Foco: NEGOCIO continúa aunque sistemas fallen",
    "Cubre: procedimientos manuales, sitio alternativo, comunicación",
    "Alcance: RRHH, finanzas, operaciones, TI",
    "Activa cuando: incidente impacta operación normal",
    "Ej. TechServ: tramitaciones manuales mientras TI recupera",
]

for i, p in enumerate(bcp_points):
    add_text(slide, p, 0.62, 1.42 + i * 0.62, 4.18, 0.55, font_size=11.5, color=D)

# Right: DRP
add_shape(slide, 5.15, 0.95, 4.5, 4.5, D)
add_shape(slide, 5.15, 0.95, 4.5, 0.06, GRN)
add_text(slide, "DRP — Plan de Recuperación ante Desastres", 5.32, 1.06, 4.18, 0.28, font_size=13, color=W, bold=True)

drp_points = [
    "Foco: SISTEMAS TIC se recuperan al estado operativo",
    "Cubre: recuperación servidores, BBDD, redes, aplicaciones",
    "Alcance: equipo TI — sysadmins, DBAs, sec ops",
    "Activa cuando: sistemas TIC fallan o son comprometidos",
    "Ej. TechServ: restaurar BBDD nóminas desde backup, RTO 8h",
]

for i, p in enumerate(drp_points):
    add_text(slide, p, 5.32, 1.42 + i * 0.62, 4.18, 0.55, font_size=11.5, color=WMUT)

# ─────────────────────────────────────────
# SLIDE 15: ENS [op.cont.*]
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "ENS Anexo II — Medidas de continuidad [op.cont.*]", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "Las 4 medidas del marco operacional de continuidad",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

measures = [
    ("[op.cont.1]", "Análisis de impacto", "Identificar funciones críticas, dependencias, RTO/RPO", "BÁSICA/MEDIA/ALTA", BLU),
    ("[op.cont.2]", "Plan de continuidad", "Procedimientos para mantener servicios críticos (BCP+DRP)", "MEDIA/ALTA", BLU),
    ("[op.cont.3]", "Pruebas periódicas", "Pruebas anuales mínimo. Documentar resultados. Actualizar.", "MEDIA/ALTA", "333333"),
    ("[op.cont.4]", "Medios alternativos", "Sitio DR, cloud, proveedores backup, o procedimientos manuales", "ALTA", "333333"),
]

for i, (code, name, req, cat, color) in enumerate(measures):
    y = 1.0 + i * 1.1
    add_shape(slide, 0.45, y, 9.1, 1.02, G if i % 2 == 0 else W)
    add_shape(slide, 0.45, y, 0.06, 1.02, color if not isinstance(color, str) else RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)))
    add_text(slide, code, 0.6, y + 0.08, 1.0, 0.22, font_size=10, color=color if not isinstance(color, str) else RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)), bold=True)
    add_text(slide, name, 0.6, y + 0.32, 1.5, 0.26, font_size=11, color=D, bold=True)
    add_text(slide, req, 2.22, y + 0.1, 5.05, 0.82, font_size=11.5, color=D)
    add_shape(slide, 7.32, y + 0.1, 0.92, 0.3, RGBColor(139, 26, 26))
    add_text(slide, cat, 7.34, y + 0.12, 0.88, 0.26, font_size=8.5, color=W, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────
# SLIDE 16: ISO 22301
# ─────────────────────────────────────────
slide = white_slide()
add_text(slide, "ISO 22301 — El estándar de gestión de continuidad", 0.5, 0.25, 9.1, 0.52,
         font_size=27, color=D, bold=True)
add_text(slide, "Cómo se relaciona con ISO 27001 y el ENS",
         0.5, 0.82, 9.1, 0.3, font_size=12.5, color=MUT)

# Left
add_shape(slide, 0.45, 0.95, 4.55, 4.42, G)
add_text(slide, "¿Qué es ISO 22301?", 0.62, 1.04, 4.22, 0.3, font_size=13, color=D, bold=True)
add_text(slide, "Estándar internacional para Gestión de Continuidad de Negocio. Define requisitos para planificar, establecer, implementar y mantener planes.\n\nEs certificable independientemente de ISO 27001.\n\nComparte estructura HLS con ISO 27001.",
         0.62, 1.42, 4.22, 3.8, font_size=12, color=D)

# Right
add_shape(slide, 5.18, 0.95, 4.42, 4.42, D)
add_text(slide, "ISO 27001 · ISO 22301 · ENS", 5.36, 1.04, 4.08, 0.3, font_size=13, color=W, bold=True)

relations = [
    ("ISO 27001", "Cubre SI incluyendo continuidad de SI (A.5.29, A.5.30)"),
    ("ISO 22301", "Cubre continuidad del negocio completo incluyendo SI"),
    ("ENS [op.cont.*]", "Requiere elementos clave (BIA, BCP, DRP, pruebas) MEDIA/ALTA"),
    ("Relación", "Son complementarios. Muchas orgs: primero ISO 27001, luego ISO 22301"),
]

for i, (t, d) in enumerate(relations):
    y = 1.42 + i * 0.82
    add_shape(slide, 5.3, y, 4.12, 0.74, DG)
    add_text(slide, t, 5.46, y + 0.08, 1.2, 0.22, font_size=10, color=TEA, bold=True)
    add_text(slide, d, 5.46, y + 0.32, 3.82, 0.38, font_size=10.5, color=WMUT)

# ─────────────────────────────────────────
# SLIDE 17: EJERCICIO BIA
# ─────────────────────────────────────────
slide = dark_slide()
add_text(slide, "Ejercicio · 8 minutos", 0.65, 0.48, 8, 0.28, font_size=12, color=SMUT, bold=True)
add_text(slide, "BIA básico para TechServ", 0.65, 0.84, 9, 0.58, font_size=34, color=W, bold=True)
add_text(slide, "Construid el análisis de impacto para los 3 sistemas principales",
         0.65, 1.5, 9, 0.3, font_size=13, color=WMUT)

add_shape(slide, 0.5, 1.92, 9.1, 3.0, DG)
add_text(slide, "Sistema | Función | Impacto 1h | Impacto 24h | RTO | RPO | Dependencias",
         0.62, 1.98, 8.83, 0.28, font_size=7.5, color=W, bold=True)

systems = ["Sistema de nóminas", "Sistema de expedientes", "Sede electrónica"]
for i, s in enumerate(systems):
    y = 2.3 + i * 0.6
    add_shape(slide, 0.62, y, 8.83, 0.54, RGBColor(45, 45, 45) if i % 2 == 0 else RGBColor(34, 34, 34))
    add_text(slide, s, 0.66, y + 0.12, 1.45, 0.28, font_size=10.5, color=W, bold=True)
    add_text(slide, "?", 2.25, y + 0.12, 0.4, 0.28, font_size=14, color=SMUT, align=PP_ALIGN.CENTER)

add_shape(slide, 0.5, 5.02, 9.1, 0.6, D)
add_text(slide, "Pista: usad RTO/RPO de análisis anterior. Para impacto pensad en consecuencias legales (ENS), operativas (servicios) y reputacionales.",
         0.7, 5.06, 8.7, 0.5, font_size=11, color=WMUT)

# ─────────────────────────────────────────
# SLIDE 18: CIERRE - 5 IDEAS
# ─────────────────────────────────────────
slide = dark_slide()
add_text(slide, "5 ideas para llevarse de la Sesión 5", 0.65, 0.35, 9, 0.52, font_size=29, color=W, bold=True)

ideas = [
    ("1", "A.8.13 sin prueba = no backup", "Backups sin restauración probada es seguridad ficticia."),
    ("2", "Etapa 1 doc, Etapa 2 evidencias", "Si Etapa 1 falla → implementación. Etapa 2 NC mayor → espera."),
    ("3", "RTO < MAD siempre", "El tiempo de recuperación < tiempo máximo tolerado. Sino, daño irreversible."),
    ("4", "BCP ≠ DRP", "DRP recupera sistemas. BCP garantiza negocio mientras TI recupera. Necesitas ambos."),
    ("5", "[op.cont.1] en TODAS categorías", "BIA obligatorio BÁSICA/MEDIA/ALTA. Sin excepción. Es el control más universal."),
]

for i, (n, t, d) in enumerate(ideas):
    if i < 3:
        x, y, w = 0.45 + i * 3.18, 1.02, 3.0
    else:
        x, y, w = 0.45 if i == 3 else 5.27, 3.0, 4.72

    add_shape(slide, x, y, w, 1.78, DG)
    add_text(slide, n, x + 0.15, y + 0.1, 0.38, 0.32, font_size=12, color=WMUT, bold=True)
    add_text(slide, t, x + 0.15, y + 0.46, w - 0.3, 0.3, font_size=12, color=W, bold=True)
    add_text(slide, d, x + 0.15, y + 0.82, w - 0.3, 0.88, font_size=10.5, color=WMUT)

add_shape(slide, 0.45, 4.87, 9.1, 0.62, D)
add_text(slide, "Próxima sesión S6 — Continuidad de Negocio completo:",
         0.65, 4.9, 3.5, 0.24, font_size=11, color=TEA, bold=True)
add_text(slide, "Taller: BCP y DRP completos · Pruebas · ISO 22301 · Gestión de crisis · Simulacro tabletop",
         4.28, 4.9, 5.1, 0.55, font_size=10.5, color=WMUT)

# ─────────────────────────────────────────
# Guardar presentación
# ─────────────────────────────────────────
output_path = Path(__file__).parent / "S5_A8_Cert_BCM.pptx"
prs.save(str(output_path))
print(f"✓ Presentación generada exitosamente: {output_path}")
print(f"✓ Total de slides: {len(prs.slides)}")
