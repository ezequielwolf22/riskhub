"""Tests de la comprension y del orquestador, con la IA mockeada.

Lo que se comprueba aqui no es que el modelo acierte (eso se valida contra la
API real), sino que el codigo que lo rodea hace lo correcto con lo que el
modelo devuelva — incluido cuando devuelve basura:

- un documento se descompone en las filas que declara su mapa,
- las unidades se vuelcan en orden de dependencia aunque el mapa venga
  desordenado,
- una regla de aplicabilidad propuesta por la IA nace DESACTIVADA,
- una entidad o una familia inventadas se descartan en vez de colarse,
- el lector se enfrenta al pack real del cliente sin romperse.
"""
from unittest.mock import patch

import pytest
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.database import Base
from app.models import (BCMApplicabilityRule, BCMLocation, BCMScenario,
                        BCMScenarioAssessment, BusinessProcess, IngestBatch,
                        IngestRecordTrace, IngestSourceMap, Supplier)
from app.services.ingest import comprehension, contracts, pipeline, reader

_ENGINE = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False})
_Session = sessionmaker(bind=_ENGINE)
Base.metadata.create_all(bind=_ENGINE)

ORG = 777001


@pytest.fixture()
def db():
    session = _Session()
    yield session
    for model in (IngestRecordTrace, IngestSourceMap, BCMScenarioAssessment,
                  BCMApplicabilityRule, BCMScenario, BusinessProcess, Supplier,
                  BCMLocation, IngestBatch):
        session.query(model).filter_by(organization_id=ORG).delete(
            synchronize_session=False)
    session.commit()
    session.close()


def _xlsx(sheets: dict) -> bytes:
    """Genera un libro real en memoria para no depender de ficheros externos."""
    import io

    import openpyxl
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── Lector ───────────────────────────────────────────────────────────────────

def test_el_lector_conserva_la_estructura_de_cada_hoja():
    data = _xlsx({
        "Escenarios": [["Escenario", "Sede", "RTO"],
                       ["Huelga", "Madrid", "4 horas"],
                       ["Caida AWS", "Madrid", "1 hora"]],
        "Proveedores": [["Proveedor", "SLA"], ["AWS", "<1h"]],
    })
    doc = reader.read_document(data, "BIA.xlsx")
    assert doc["format"] == "xlsx"
    assert doc["stats"]["sheets"] == 2
    rendered = reader.render_for_llm(doc)
    # La hoja, su referencia y el numero de fila real llegan al modelo: son lo
    # que le permite decir "las filas 2 y 3 son dos unidades".
    assert "sheet:Escenarios" in rendered
    assert "f2: Huelga | Madrid | 4 horas" in rendered
    assert "sheet:Proveedores" in rendered


def test_los_errores_de_formula_no_son_datos():
    data = _xlsx({"H": [["A", "B"], ["#REF!", "valido"]]})
    rendered = reader.render_for_llm(reader.read_document(data, "x.xlsx"))
    assert "#REF!" not in rendered
    assert "valido" in rendered


def test_formato_no_soportado_se_declara_en_vez_de_romper():
    assert reader.is_supported("diagrama.vsdx") is False
    with pytest.raises(reader.UnsupportedDocument):
        reader.read_document(b"x", "diagrama.vsdx")


def test_lee_presentaciones():
    """Los packs de continuidad traen el flujo de crisis como presentacion."""
    import io as _io

    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "WorkFlow Indisponibilidad de personas"
    slide.placeholders[1].text = "Activar el Plan de Comunicacion de crisis"
    buf = _io.BytesIO()
    prs.save(buf)

    doc = reader.read_document(buf.getvalue(), "workflow.pptx")
    assert doc["format"] == "pptx"
    rendered = reader.render_for_llm(doc)
    assert "slide:1" in rendered
    assert "WorkFlow Indisponibilidad de personas" in rendered
    assert "Plan de Comunicacion de crisis" in rendered


def test_el_catalogo_del_prompt_sale_del_registro():
    described = contracts.describe_for_prompt()
    assert "bcm_scenario_assessment" in described
    # Los campos calculados se declaran como no extraibles
    assert "NO extraigas" in described
    assert "weighted_impact" in described
    # Las referencias se piden por nombre, nunca por id
    assert "_ref_scenario_id" in described


# ── Saneado de lo que devuelve el modelo ─────────────────────────────────────

def test_una_entidad_inventada_se_descarta():
    parsed = {"units": [
        {"target_entity": "tabla_que_no_existe", "label": "x",
         "decomposition_key": "y", "rows": [{"fields": {"name": "A"}}]},
        {"target_entity": "business_process", "label": "ok",
         "decomposition_key": "y", "rows": [{"fields": {"name": "B"}}]},
    ]}
    out = comprehension._sanitize_map(parsed, {"filename": "f", "sha256": "s"})
    assert len(out["units"]) == 1
    assert out["units"][0]["target_entity"] == "business_process"


def test_una_regla_mal_formada_no_llega_a_guardarse():
    parsed = {"applicability_rules": [
        {"when": {"planta": ["baja"]}, "then": {"exclude_families": ["facilities"]},
         "rationale": "campo de condicion inventado"},
        {"when": {"site_type": ["remote"]}, "then": {"exclude_families": ["cosmico"]},
         "rationale": "familia inventada"},
        {"when": {"site_type": ["remote"]}, "then": {},
         "rationale": "sin efecto"},
        {"when": {"site_type": ["remote"]},
         "then": {"exclude_families": ["facilities"]}, "rationale": ""},
        {"when": {"site_type": ["remote"]},
         "then": {"include_only_families": ["personnel"]},
         "rationale": "Sin oficina propia no hay escenario de instalaciones."},
    ]}
    out = comprehension._sanitize_profile(parsed)
    assert len(out["applicability_rules"]) == 1
    assert out["applicability_rules"][0]["then"] == {"include_only_families": ["personnel"]}


# ── Orquestador ──────────────────────────────────────────────────────────────

_PROFILE = {
    "narrative": "Grupo con sedes en Espana y Latam.",
    "confidence": 0.9,
    "structure": {"locations": [
        {"name": "Madrid - HQ", "country": "Espana", "city": "Madrid",
         "site_type": "office", "evidence": "portada del BIA"},
        {"name": "Santiago", "country": "Chile", "site_type": "remote"},
    ]},
    "bia_method": {"style": "scenario_x_location",
                   "horizons": ["0h", ">4h"],
                   "aggregation": "max"},
    "scenarios": [
        {"code": "ALT.01", "name": "Huelga de personal", "family": "personnel"},
        {"code": "ALT.07", "name": "Caida de la infraestructura AWS",
         "family": "systems_comms"},
    ],
    "applicability_rules": [
        {"when": {"site_type": ["remote"]},
         "then": {"include_only_families": ["personnel"]},
         "rationale": "Las sedes remotas no tienen instalaciones propias."},
    ],
}

# El mapa llega deliberadamente en orden inverso al de dependencia: primero las
# valoraciones y despues los escenarios que referencian.
_MAP = {
    "doc_kind": "bia",
    "confidence": 0.9,
    "rationale": "Cada bloque de la hoja es un escenario valorado en una sede.",
    "filename": "BIA.xlsx",
    "sha256": "abc123",
    "units": [
        {"label": "Valoraciones", "source_ref": "sheet:Escenarios",
         "decomposition_key": "un bloque por escenario, clave en la columna B",
         "target_entity": "bcm_scenario_assessment", "confidence": 0.9,
         "rows": [
             {"fields": {"_ref_scenario_id": "ALT.01",
                         "_ref_location_id": "Oficina de Madrid",
                         "rto_label": "4 horas", "rto_hours": 4,
                         "impacts": {"operational": {">4h": 5}}},
              "source_ref": "sheet:Escenarios!f8"},
             {"fields": {"_ref_scenario_id": "ALT.07",
                         "_ref_location_id": "Madrid - HQ",
                         "rto_label": "1 hora", "rto_hours": 1,
                         "impacts": {"operational": {">4h": 4}}},
              "confidence": 0.5},
         ]},
        {"label": "Escenarios", "source_ref": "sheet:Escenarios",
         "decomposition_key": "una fila por escenario",
         "target_entity": "bcm_scenario", "confidence": 0.95,
         "rows": [
             {"fields": {"code": "ALT.01", "name": "Huelga de personal",
                         "family": "personnel"}},
             {"fields": {"code": "ALT.07", "name": "Caida de la infraestructura AWS",
                         "family": "systems_comms"}},
         ]},
    ],
    "ambiguities": [{"question": "La columna E es RTO o MTPD?",
                     "chosen": "RTO", "why": "la cabecera dice RTO"}],
}


def _run(db, **kw):
    data = _xlsx({"Escenarios": [["Escenario", "Sede"], ["Huelga", "Madrid"]]})
    with patch.object(comprehension, "build_profile", return_value=_PROFILE), \
         patch.object(comprehension, "build_source_map", return_value=_MAP):
        return pipeline.run_pack(db, ORG, [("BIA.xlsx", data)], **kw)


def test_el_pack_se_vuelca_completo_y_en_orden_de_dependencia(db):
    out = _run(db)
    assert out["status"] == "completed"

    # Las sedes y los escenarios del perfil existen
    assert db.query(BCMLocation).filter_by(organization_id=ORG).count() == 2
    assert db.query(BCMScenario).filter_by(organization_id=ORG).count() == 2

    # Y las valoraciones encontraron su escenario y su sede pese a que el mapa
    # venia en orden inverso y la sede se nombraba de otra forma
    valoraciones = db.query(BCMScenarioAssessment).filter_by(
        organization_id=ORG).all()
    assert len(valoraciones) == 2
    madrid = db.query(BCMLocation).filter_by(
        organization_id=ORG, name="Madrid - HQ").one()
    assert all(v.scenario_id and v.location_id == madrid.id for v in valoraciones)


def test_el_impacto_ponderado_lo_pone_el_motor_tras_volcar(db):
    _run(db)
    v = db.query(BCMScenarioAssessment).filter_by(organization_id=ORG).filter(
        BCMScenarioAssessment.rto_label == "4 horas").one()
    # nivel 5 -> score 4, factor del RTO "4 horas" = 0.6
    assert v.weighted_impact == pytest.approx(2.4)
    assert v.impact_band == "severe"


def test_una_fila_dudosa_se_guarda_marcada_no_se_pierde(db):
    out = _run(db)
    assert out["needs_review"] >= 1
    dudosa = db.query(BCMScenarioAssessment).filter_by(
        organization_id=ORG, needs_review=True).one()
    assert dudosa.import_confidence == 0.5


def test_la_regla_propuesta_por_la_ia_nace_desactivada(db):
    """Una regla activa oculta escenarios: eso no lo decide el agente solo."""
    out = _run(db)
    rule = db.query(BCMApplicabilityRule).filter_by(organization_id=ORG).one()
    assert rule.is_active is False
    assert rule.source == "proposed"
    assert rule.rationale
    assert any("desactivada" in w for w in out["warnings"])


def test_se_guarda_el_como_lo_entendi_con_su_razonamiento(db):
    out = _run(db)
    smap = db.query(IngestSourceMap).filter_by(organization_id=ORG).one()
    assert smap.doc_kind == "bia"
    assert smap.rationale
    assert smap.ambiguities
    assert {u["target_entity"] for u in smap.units_detected} == {
        "bcm_scenario", "bcm_scenario_assessment"}
    assert out["batch_id"]


def test_la_verificacion_declara_los_huecos_que_quedan(db):
    out = _run(db)
    gaps = out["gaps"]
    # Dos escenarios sin estrategia, sin plan y sin ejercitar en la sede con
    # oficina; la remota queda cubierta por la regla... que esta desactivada,
    # asi que tambien cuenta. Los huecos se declaran, no se maquillan.
    assert gaps["total"] > 0
    assert "no_strategy" in gaps["by_reason"]
    assert "not_tested_12m" in gaps["by_reason"]


def test_sin_perfil_el_volcado_sigue_funcionando(db):
    """apply_profile=False: el cliente ya tiene su estructura montada."""
    out = _run(db, apply_profile=False)
    assert out["status"] == "completed"
    assert db.query(BCMLocation).filter_by(organization_id=ORG).count() == 0
    # Los escenarios del propio documento si se crean
    assert db.query(BCMScenario).filter_by(organization_id=ORG).count() == 2


def test_todo_el_pack_es_reversible(db):
    from app.services.ingest.batch import undo_batch
    out = _run(db)
    bat = db.get(IngestBatch, out["batch_id"])
    result = undo_batch(db, bat)
    assert result["reverted"] > 0
    assert db.query(BCMScenario).filter_by(organization_id=ORG).count() == 0
    assert db.query(BCMScenarioAssessment).filter_by(organization_id=ORG).count() == 0
    assert db.query(BCMLocation).filter_by(organization_id=ORG).count() == 0


def test_un_documento_ilegible_no_tumba_el_pack(db):
    data = _xlsx({"H": [["A"], ["1"]]})
    with patch.object(comprehension, "build_profile", return_value=_PROFILE), \
         patch.object(comprehension, "build_source_map", return_value=_MAP):
        out = pipeline.run_pack(db, ORG, [
            ("roto.xlsx", b"esto no es un xlsx"),
            ("BIA.xlsx", data),
        ])
    assert out["status"] == "completed"
    assert any("roto.xlsx" in w for w in out["warnings"])
    assert db.query(BCMScenario).filter_by(organization_id=ORG).count() == 2


def test_si_la_ia_falla_el_lote_lo_registra_sin_reventar(db):
    data = _xlsx({"H": [["A"], ["1"]]})
    with patch.object(comprehension, "build_profile",
                      side_effect=RuntimeError("sin creditos")), \
         patch.object(comprehension, "build_source_map",
                      side_effect=RuntimeError("sin creditos")):
        out = pipeline.run_pack(db, ORG, [("BIA.xlsx", data)])
    assert out["status"] == "completed"
    assert any("sin creditos" in w for w in out["warnings"])
    bat = db.get(IngestBatch, out["batch_id"])
    assert bat.files[0]["status"] == "failed"


# ── Estimacion previa ────────────────────────────────────────────────────────

def test_la_estimacion_no_llama_a_la_ia_y_declara_lo_no_soportado():
    data = _xlsx({"H": [["A", "B"], ["1", "2"]]})
    est = pipeline.estimate_pack([("BIA.xlsx", data), ("diagrama.vsdx", b"x")])
    assert est["documents_total"] == 1
    assert est["unsupported"] == ["diagrama.vsdx"]
    assert est["total_chars"] > 0
    assert est["estimated_cost_usd"] >= 0
