"""Tests del nucleo generico de ingesta.

Fijan las cuatro garantias que el usuario tiene sobre todo lo que escriba el
agente a partir de un documento:

1. Un documento se descompone en N filas, no en una.
2. Reconciliar no duplica: "AMAZON WEB SERVICES, INC." es el proveedor que ya
   existe, no uno nuevo.
3. Cuando dos documentos se contradicen gana el mas restrictivo, y el valor
   descartado queda registrado con su fuente.
4. Todo se puede deshacer, y una correccion manual sobrevive a reimportar.

BD SQLite en memoria propia, al estilo de test_bcp_scoring.py.
"""
import pytest
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.database import Base
from app.models import (BCMLocation, BCMScenario, BCMScenarioAssessment, BCPPlan,
                        BCPSupplierLink, BusinessProcess, IngestBatch,
                        IngestConflict, IngestFieldOverride, IngestRecordTrace,
                        Supplier)
from app.services.ingest import batch as batch_mod
from app.services.ingest import conflicts as conflicts_mod
from app.services.ingest import contracts
from app.services.ingest.materializer import MaterializationResult, materialize
from app.services.ingest.reconciler import find_match, normalize_name, similarity

_ENGINE = create_engine("sqlite:///:memory:", connect_args={"check_same_thread": False})
_Session = sessionmaker(bind=_ENGINE)
Base.metadata.create_all(bind=_ENGINE)

ORG = 555001
OTHER_ORG = 555002


@pytest.fixture()
def db():
    session = _Session()
    yield session
    from app.models import AiDecisionSignal
    for model in (AiDecisionSignal, IngestRecordTrace, IngestConflict, IngestFieldOverride,
                  BCMScenarioAssessment, BCMScenario, BCPSupplierLink,
                  BCPPlan, BusinessProcess, Supplier, BCMLocation, IngestBatch):
        session.query(model).filter(
            model.organization_id.in_([ORG, OTHER_ORG])
        ).delete(synchronize_session=False)
    session.commit()
    session.close()


@pytest.fixture()
def bat(db):
    return batch_mod.create_batch(db, ORG, module="bcm", files=[])


# ── Normalizacion y similitud ────────────────────────────────────────────────

def test_normalizacion_quita_el_ruido_juridico():
    assert normalize_name("AMAZON WEB SERVICES, INC.") == "amazon web services"
    assert normalize_name("TELEFONICA DE ESPAÑA S.A.U.") == "telefonica de espana"
    assert normalize_name("Dinahosting S.L.") == "dinahosting"
    assert normalize_name("  Google   Cloud  EMEA  Limited ") == "google cloud emea"


def test_similitud_reconoce_contencion_pero_no_casualidades():
    assert similarity("amazon web services", "amazon web services") == 1.0
    assert similarity("amazon web services", "amazon web services europe") > 0.8
    assert similarity("nalanda", "dokify") < 0.5


# ── Descomposicion: un documento, muchas filas ───────────────────────────────

def test_un_documento_genera_muchas_filas(db, bat):
    """El caso central: un BIA con doce procesos son doce filas, no una."""
    rows = [
        {"name": f"Proceso {i}", "criticality": "high", "rto_hours": i,
         "_source_ref": f"sheet:BIA!f{i}"}
        for i in range(1, 13)
    ]
    result = materialize(db, ORG, bat, "business_process", rows,
                         source_filename="BIA.xlsx")
    db.commit()

    assert result.created["business_process"] == 12
    assert db.query(BusinessProcess).filter_by(organization_id=ORG).count() == 12
    # Cada fila deja su propio rastro reversible
    assert db.query(IngestRecordTrace).filter_by(batch_id=bat.id).count() == 12


def test_reimportar_el_mismo_documento_no_duplica(db, bat):
    rows = [{"name": "Facturacion", "criticality": "critical", "rto_hours": 4}]
    materialize(db, ORG, bat, "business_process", rows, source_filename="BIA.xlsx")
    db.commit()
    materialize(db, ORG, bat, "business_process", rows, source_filename="BIA.xlsx")
    db.commit()
    assert db.query(BusinessProcess).filter_by(organization_id=ORG).count() == 1


# ── Reconciliacion ───────────────────────────────────────────────────────────

def test_reconcilia_el_proveedor_existente_en_vez_de_duplicarlo(db, bat):
    db.add(Supplier(organization_id=ORG, code="SUP-0001", name="Amazon Web Services"))
    db.commit()

    result = materialize(db, ORG, bat, "supplier",
                         [{"name": "AMAZON WEB SERVICES, INC.",
                           "services": "infraestructura"}],
                         source_filename="BIA.xlsx")
    db.commit()

    assert db.query(Supplier).filter_by(organization_id=ORG).count() == 1
    assert result.created.get("supplier") is None
    sup = db.query(Supplier).filter_by(organization_id=ORG).one()
    # Enlazar no pierde el dato nuevo: el campo vacio se rellena
    assert sup.services == "infraestructura"


def test_la_reconciliacion_no_cruza_organizaciones(db, bat):
    db.add(Supplier(organization_id=OTHER_ORG, code="SUP-0002", name="Amazon Web Services"))
    db.commit()

    materialize(db, ORG, bat, "supplier", [{"name": "AMAZON WEB SERVICES, INC."}])
    db.commit()

    assert db.query(Supplier).filter_by(organization_id=ORG).count() == 1
    assert db.query(Supplier).filter_by(organization_id=OTHER_ORG).count() == 1


def test_un_parecido_nunca_se_fusiona_solo(db, bat):
    """Un nombre parecido no es prueba de ser el mismo registro.

    Antes, un parecido >= 0.88 auto-actualizaba el registro existente y
    aplastaba datos. Ahora nunca se fusiona por parecido: se marca como posible
    duplicado y la decision es del usuario.
    """
    db.add_all([Supplier(organization_id=ORG, code="SUP-0003", name="Grupo Solana"),
                Supplier(organization_id=ORG, code="SUP-0004", name="Grupo Solano")])
    db.commit()

    spec = contracts.get("supplier")
    match = find_match(db, spec, ORG, {"name": "Grupo Solane"})
    assert match.how == "possible_duplicate"
    assert match.matched is False           # no se enlaza a ciegas
    assert len(match.candidates) >= 1


def test_un_nombre_identico_normalizado_si_enlaza(db, bat):
    """El caso legitimo: 'AMAZON WEB SERVICES, INC.' == 'Amazon Web Services'."""
    db.add(Supplier(organization_id=ORG, code="SUP-0009", name="Amazon Web Services"))
    db.commit()
    match = find_match(db, contracts.get("supplier"), ORG,
                       {"name": "AMAZON WEB SERVICES, INC."})
    assert match.matched is True
    assert match.how == "natural_key"


def test_referencias_por_clave_natural_se_resuelven_dentro_del_lote(db, bat):
    """El documento dice "escenario ALT.05"; en la base va el id."""
    materialize(db, ORG, bat, "bcm_scenario",
                [{"code": "ALT.05", "name": "Caida de las comunicaciones",
                  "family": "systems_comms"}])
    materialize(db, ORG, bat, "bcm_location",
                [{"name": "Madrid - HQ", "country": "Espana", "site_type": "office"}])
    result = materialize(db, ORG, bat, "bcm_scenario_assessment",
                         [{"_ref_scenario_id": "ALT.05",
                           "_ref_location_id": "Oficina de Madrid",
                           "rto_label": "4 horas", "rto_hours": 4,
                           "impacts": {"operational": {">6h": 4}}}])
    db.commit()

    assert result.created.get("bcm_scenario_assessment") == 1
    a = db.query(BCMScenarioAssessment).filter_by(organization_id=ORG).one()
    sc = db.query(BCMScenario).filter_by(organization_id=ORG).one()
    loc = db.query(BCMLocation).filter_by(organization_id=ORG).one()
    assert a.scenario_id == sc.id
    # "Oficina de Madrid" y "Madrid - HQ" son la misma sede
    assert a.location_id == loc.id


# ── Campos calculados y vocabularios ─────────────────────────────────────────

def test_un_campo_calculado_no_se_escribe_desde_el_documento(db, bat):
    materialize(db, ORG, bat, "bcm_scenario",
                [{"code": "ALT.01", "name": "Huelga", "family": "personnel"}])
    materialize(db, ORG, bat, "bcm_scenario_assessment",
                [{"_ref_scenario_id": "ALT.01",
                  "impacts": {"operational": {">6h": 5}},
                  "weighted_impact": 99.0, "impact_band": "trivial"}])
    db.commit()

    a = db.query(BCMScenarioAssessment).filter_by(organization_id=ORG).one()
    assert a.weighted_impact is None
    assert a.impact_band is None


def test_valor_fuera_del_vocabulario_se_descarta(db, bat):
    materialize(db, ORG, bat, "business_process",
                [{"name": "Proceso raro", "criticality": "importantisimo"}])
    db.commit()
    p = db.query(BusinessProcess).filter_by(organization_id=ORG).one()
    assert p.criticality != "importantisimo"


def test_el_estado_de_un_plan_nunca_se_importa(db, bat):
    """Que un documento se titule "aprobado" no aprueba nada en RiskHub."""
    materialize(db, ORG, bat, "bcp_plan",
                [{"name": "DRP AWS", "plan_type": "drp", "status": "approved"}])
    db.commit()
    plan = db.query(BCPPlan).filter_by(organization_id=ORG).one()
    assert plan.status == "draft"


def test_confianza_baja_marca_para_revision_pero_escribe(db, bat):
    result = materialize(db, ORG, bat, "business_process",
                         [{"name": "Proceso dudoso", "_confidence": 0.4}])
    db.commit()
    assert result.created["business_process"] == 1
    assert result.needs_review == 1
    p = db.query(BusinessProcess).filter_by(organization_id=ORG).one()
    assert p.needs_review is True
    assert p.import_confidence == 0.4


# ── Conflictos ───────────────────────────────────────────────────────────────

def test_en_un_rto_gana_el_mas_exigente():
    c = [conflicts_mod.Candidate(6, source_filename="DRP.docx"),
         conflicts_mod.Candidate(4, source_filename="BIA.xlsx")]
    res = conflicts_mod.resolve("rto_hours", "min", c)
    assert res["value"] == 4
    assert res["decided"] is True
    assert len(res["discarded"]) == 1
    assert res["discarded"][0]["source_filename"] == "DRP.docx"


def test_en_una_criticidad_gana_la_mas_alta():
    c = [conflicts_mod.Candidate("medium"), conflicts_mod.Candidate("critical")]
    assert conflicts_mod.resolve("criticality", "max", c)["value"] == "critical"


def test_politica_manual_no_decide():
    c = [conflicts_mod.Candidate("A"), conflicts_mod.Candidate("B")]
    res = conflicts_mod.resolve("x", "manual", c)
    assert res["decided"] is False and res["value"] is None


def test_fuentes_que_coinciden_no_son_conflicto():
    c = [conflicts_mod.Candidate(4), conflicts_mod.Candidate(4)]
    assert conflicts_mod.resolve("rto_hours", "min", c)["discarded"] == []


def test_documentos_contradictorios_dejan_el_conflicto_registrado(db, bat):
    """BIA dice 4h, DRP dice 6h: gana 4 y el 6 no se pierde."""
    materialize(db, ORG, bat, "business_process",
                [{"name": "Plataforma CAE", "rto_hours": 4}],
                source_filename="BIA.xlsx")
    db.commit()
    materialize(db, ORG, bat, "business_process",
                [{"name": "Plataforma CAE", "rto_hours": 6}],
                source_filename="DRP.docx")
    db.commit()

    p = db.query(BusinessProcess).filter_by(organization_id=ORG).one()
    assert p.rto_hours == 4
    assert p.needs_review is True

    conflict = db.query(IngestConflict).filter_by(
        organization_id=ORG, field_name="rto_hours").one()
    assert conflict.resolved_value == 4
    assert conflict.policy == "min"
    fuentes = {c.get("source_filename") for c in conflict.candidates}
    assert "DRP.docx" in fuentes


def test_el_conflicto_cita_el_documento_real_del_valor_existente(db, bat):
    """Subiendo todo de golpe, el valor 'ya cargado' lo puso otro documento del
    pack: el conflicto debe citar ESE documento, no un generico."""
    # El pipeline comparte un mismo result entre documentos: aqui igual.
    result = MaterializationResult()
    materialize(db, ORG, bat, "business_process",
                [{"name": "Plataforma CAE", "rto_hours": 4}],
                source_filename="BIA_Nalanda.xlsx", result=result)
    db.commit()
    materialize(db, ORG, bat, "business_process",
                [{"name": "Plataforma CAE", "rto_hours": 6}],
                source_filename="ISRT_03.docx", result=result)
    db.commit()

    conflict = db.query(IngestConflict).filter_by(
        organization_id=ORG, field_name="rto_hours").one()
    by_value = {c["value"]: c for c in conflict.candidates}
    # El valor que ya estaba (4) lleva el documento que lo puso, no "valor ya cargado"
    assert by_value[4]["source_filename"] == "BIA_Nalanda.xlsx"
    assert by_value[6]["source_filename"] == "ISRT_03.docx"


def test_valor_preexistente_sin_rastro_no_finge_un_documento(db, bat):
    """Si el valor venia de una importacion anterior (sin procedencia en este
    pack), el candidato va sin documento; la UI lo etiqueta aparte."""
    db.add(BusinessProcess(organization_id=ORG, name="Portal CAE", rto_hours=8))
    db.commit()
    materialize(db, ORG, bat, "business_process",
                [{"name": "Portal CAE", "rto_hours": 2}],
                source_filename="DRP.docx", result=MaterializationResult())
    db.commit()

    conflict = db.query(IngestConflict).filter_by(
        organization_id=ORG, field_name="rto_hours").one()
    by_value = {c["value"]: c for c in conflict.candidates}
    assert by_value[8]["source_filename"] is None       # sin fingir documento
    assert by_value[2]["source_filename"] == "DRP.docx"


# ── Deshacer y forzar ────────────────────────────────────────────────────────

def test_deshacer_el_lote_deja_la_base_como_estaba(db, bat):
    db.add(Supplier(organization_id=ORG, code="SUP-0005", name="Proveedor previo"))
    db.commit()

    materialize(db, ORG, bat, "business_process",
                [{"name": "P1"}, {"name": "P2"}, {"name": "P3"}])
    materialize(db, ORG, bat, "supplier", [{"name": "Proveedor nuevo"}])
    db.commit()
    assert db.query(BusinessProcess).filter_by(organization_id=ORG).count() == 3

    res = batch_mod.undo_batch(db, bat)
    assert res["reverted"] == 4
    assert res["failed"] == []
    assert db.query(BusinessProcess).filter_by(organization_id=ORG).count() == 0
    # Lo que ya existia antes del lote sigue intacto
    assert db.query(Supplier).filter_by(organization_id=ORG).count() == 1
    assert bat.status == "undone"


def test_deshacer_restaura_el_valor_previo_de_lo_actualizado(db, bat):
    # Un registro con un hueco (rto sin informar): la ingesta RELLENA el hueco,
    # que es la unica forma en que una importacion cambia un registro existente.
    db.add(BusinessProcess(organization_id=ORG, name="Facturacion"))
    db.commit()

    materialize(db, ORG, bat, "business_process",
                [{"name": "Facturacion", "rto_hours": 4}])
    db.commit()
    assert db.query(BusinessProcess).filter_by(organization_id=ORG).one().rto_hours == 4

    batch_mod.undo_batch(db, bat)
    p = db.query(BusinessProcess).filter_by(organization_id=ORG).one()
    assert p.rto_hours is None   # vuelve al hueco original, no se inventa nada


def test_revertir_un_solo_registro(db, bat):
    materialize(db, ORG, bat, "business_process", [{"name": "A"}, {"name": "B"}])
    db.commit()
    trace = db.query(IngestRecordTrace).filter_by(batch_id=bat.id).first()
    assert batch_mod.revert_record(db, trace) is True
    db.commit()
    assert db.query(BusinessProcess).filter_by(organization_id=ORG).count() == 1


def test_una_correccion_manual_sobrevive_a_reimportar(db, bat):
    """La garantia de "puedo forzar cambios y no se pierden"."""
    materialize(db, ORG, bat, "business_process",
                [{"name": "Nominas", "rto_hours": 24}], source_filename="BIA.xlsx")
    db.commit()
    p = db.query(BusinessProcess).filter_by(organization_id=ORG).one()

    # El usuario corrige: sabe que en realidad son 8 horas
    p.rto_hours = 8
    batch_mod.set_override(db, ORG, "business_processes", p.id, "rto_hours",
                           value=8, previous_value=24,
                           reason="Acordado con el responsable del proceso")
    db.commit()

    # Se vuelve a importar el mismo documento, que sigue diciendo 24
    bat2 = batch_mod.create_batch(db, ORG, module="bcm")
    materialize(db, ORG, bat2, "business_process",
                [{"name": "Nominas", "rto_hours": 24}], source_filename="BIA.xlsx")
    db.commit()

    db.refresh(p)
    assert p.rto_hours == 8
    # Y no se registra como conflicto: la decision ya estaba tomada
    assert db.query(IngestConflict).filter_by(
        organization_id=ORG, field_name="rto_hours").count() == 0


def test_forzar_un_campo_emite_senal_de_aprendizaje(db, bat):
    from app.models import AiDecisionSignal
    materialize(db, ORG, bat, "business_process", [{"name": "Tesoreria"}])
    db.commit()
    p = db.query(BusinessProcess).filter_by(organization_id=ORG).one()

    batch_mod.set_override(db, ORG, "business_processes", p.id, "criticality",
                           value="critical", previous_value="medium",
                           reason="Toda tesoreria es critica en esta casa")
    db.commit()

    signal = db.query(AiDecisionSignal).filter_by(
        organization_id=ORG, signal_type="ingest_field_override").one()
    assert signal.context["user_value"] == "critical"
    assert signal.context["field"] == "criticality"


# ── Contrato ─────────────────────────────────────────────────────────────────

def test_las_entidades_bcp_estan_registradas():
    specs = contracts.all_specs()
    for key in ("bcm_location", "bcm_scenario", "bcm_scenario_assessment",
                "business_process", "bcp_dependency", "bcp_strategy",
                "bcp_plan", "bcp_test", "supplier", "bcp_supplier_link"):
        assert key in specs, f"falta la entidad {key}"


def test_las_dependencias_se_ordenan_antes_que_quien_las_usa():
    order = contracts.ordered_keys()
    assert order.index("bcm_scenario") < order.index("bcm_scenario_assessment")
    assert order.index("business_process") < order.index("bcp_dependency")
    assert order.index("supplier") < order.index("bcp_supplier_link")


def test_los_objetivos_de_recuperacion_ganan_por_el_mas_exigente():
    """Regla de dominio, no detalle: en continuidad se peca de prudente."""
    for key, fields in (("business_process", ("rto_hours", "rpo_hours", "mtpd_hours")),
                        ("bcm_scenario_assessment", ("rto_hours", "rpo_hours"))):
        fmap = contracts.get(key).field_map()
        for fname in fields:
            assert fmap[fname].conflict_policy == "min", f"{key}.{fname}"


# ── Lectura de documentos ────────────────────────────────────────────────────

def test_pptx_con_tabla_se_lee_sin_romperse():
    """Regresion: `_RowCollection` de python-pptx no soporta slicing.

    Los packs de crisis suelen traer medios alternativos y escalado como
    presentacion con tablas; antes cualquier PPTX con tabla tumbaba el lector.
    """
    import io

    from pptx import Presentation
    from pptx.util import Inches

    from app.services.ingest import reader

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Escalado de crisis"
    s1.placeholders[1].text = "Nivel alto: Comite de crisis"
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Medios alternativos"
    tbl = s2.shapes.add_table(2, 3, Inches(0.5), Inches(1.5),
                              Inches(9), Inches(1.5)).table
    for j, h in enumerate(["Medio", "Transicion", "Tiempo"]):
        tbl.cell(0, j).text = h
    tbl.cell(1, 0).text = "Trabajo remoto"
    tbl.cell(1, 1).text = "Manual"
    tbl.cell(1, 2).text = "120 minutos"
    buf = io.BytesIO()
    prs.save(buf)

    doc = reader.read_document(buf.getvalue(), "crisis.pptx")
    assert doc["format"] == "pptx"
    tables = [b for b in doc["blocks"] if b.get("type") == "table"]
    assert tables, "la tabla de la diapositiva deberia extraerse"
    assert tables[0]["header"] == ["Medio", "Transicion", "Tiempo"]
    assert ["Trabajo remoto", "Manual", "120 minutos"] in tables[0]["rows"]
